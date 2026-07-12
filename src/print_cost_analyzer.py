"""Análisis técnico y costeo estimado de archivos para impresión.

El cálculo de cobertura es físico para imágenes y PDF renderizado. En DOCX/XLSX/PPTX
se usa una estimación basada en imágenes y elementos del documento cuando no existe
un motor de renderizado de Office disponible.
"""
from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Iterable
import math
import zipfile

import streamlit as st
from PIL import Image, ImageStat


MAX_FILE_MB = 40
SUPPORTED = ["pdf", "jpg", "jpeg", "png", "docx", "xlsx", "pptx"]


@dataclass
class Coverage:
    pages: int
    cyan: float
    magenta: float
    yellow: float
    black: float
    confidence: str = "Alta"
    note: str = ""

    @property
    def total(self) -> float:
        return self.cyan + self.magenta + self.yellow + self.black


def _image_coverage(image: Image.Image) -> tuple[float, float, float, float]:
    rgb = image.convert("RGB")
    rgb.thumbnail((1400, 1400))
    cmyk = rgb.convert("CMYK")
    means = ImageStat.Stat(cmyk).mean
    return tuple(max(0.0, min(100.0, value / 255 * 100)) for value in means)  # type: ignore[return-value]


def _average(values: Iterable[tuple[float, float, float, float]]) -> tuple[float, float, float, float]:
    rows = list(values)
    if not rows:
        return 0.0, 0.0, 0.0, 0.0
    return tuple(sum(row[i] for row in rows) / len(rows) for i in range(4))  # type: ignore[return-value]


def _analyze_pdf(data: bytes) -> Coverage:
    try:
        import fitz  # PyMuPDF
    except ImportError as exc:
        raise RuntimeError("Falta PyMuPDF para analizar PDF.") from exc
    doc = fitz.open(stream=data, filetype="pdf")
    covers = []
    for page in doc:
        pix = page.get_pixmap(matrix=fitz.Matrix(1.35, 1.35), alpha=False)
        image = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
        covers.append(_image_coverage(image))
    c, m, y, k = _average(covers)
    return Coverage(len(doc), c, m, y, k, "Alta", "PDF renderizado página por página.")


def _zip_images(data: bytes) -> list[Image.Image]:
    images: list[Image.Image] = []
    with zipfile.ZipFile(BytesIO(data)) as archive:
        for name in archive.namelist():
            lower = name.lower()
            if "/media/" not in lower or not lower.endswith((".png", ".jpg", ".jpeg")):
                continue
            try:
                images.append(Image.open(BytesIO(archive.read(name))).convert("RGB"))
            except Exception:
                continue
    return images


def _analyze_docx(data: bytes) -> Coverage:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("Falta python-docx para analizar Word.") from exc
    doc = Document(BytesIO(data))
    text_chars = sum(len(p.text) for p in doc.paragraphs)
    pages = max(1, math.ceil(text_chars / 2800))
    images = _zip_images(data)
    image_cov = _average(_image_coverage(img) for img in images)
    image_weight = min(0.75, len(images) / max(1, pages) * 0.22)
    text_k = min(18.0, 2.5 + text_chars / max(1, pages) / 220)
    c, m, y, k = (v * image_weight for v in image_cov)
    return Coverage(pages, c, m, y, min(100.0, k + text_k), "Media", "Word se estima por texto e imágenes incrustadas; no sustituye una previsualización renderizada.")


def _analyze_xlsx(data: bytes) -> Coverage:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError("Falta openpyxl para analizar Excel.") from exc
    wb = load_workbook(BytesIO(data), read_only=False, data_only=True)
    pages = 0
    colored = 0
    populated = 0
    for ws in wb.worksheets:
        used = max(1, ws.max_row * ws.max_column)
        populated_sheet = 0
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is not None:
                    populated += 1
                    populated_sheet += 1
                    fill = getattr(cell.fill, "fgColor", None)
                    if fill and getattr(fill, "rgb", None) not in (None, "00000000", "00FFFFFF", "FFFFFFFF"):
                        colored += 1
        pages += max(1, math.ceil(populated_sheet / 55))
    ratio = colored / max(1, populated)
    k = min(30.0, 4.0 + populated / max(1, pages) / 18)
    color = min(24.0, ratio * 45)
    return Coverage(max(1, pages), color * .85, color, color * .65, k, "Media-baja", "Excel se estima por celdas usadas, texto y rellenos; áreas de impresión y gráficos complejos pueden variar.")


def _analyze_pptx(data: bytes) -> Coverage:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("Falta python-pptx para analizar PowerPoint.") from exc
    prs = Presentation(BytesIO(data))
    images = _zip_images(data)
    image_cov = _average(_image_coverage(img) for img in images)
    text_chars = 0
    shapes = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            shapes += 1
            if hasattr(shape, "text"):
                text_chars += len(shape.text or "")
    density = min(1.0, (len(images) * 1.6 + shapes * .08) / max(1, len(prs.slides)))
    c, m, y, k = (min(100.0, v * max(.25, density)) for v in image_cov)
    k = min(100.0, k + min(18.0, text_chars / max(1, len(prs.slides)) / 90))
    return Coverage(max(1, len(prs.slides)), c, m, y, k, "Media", "PowerPoint se estima por diapositivas, imágenes y texto; degradados y vectores pueden variar.")


def analyze_file(name: str, data: bytes) -> Coverage:
    suffix = Path(name).suffix.lower()
    if suffix in (".jpg", ".jpeg", ".png"):
        c, m, y, k = _image_coverage(Image.open(BytesIO(data)))
        return Coverage(1, c, m, y, k, "Alta", "Imagen raster analizada píxel por píxel.")
    if suffix == ".pdf":
        return _analyze_pdf(data)
    if suffix == ".docx":
        return _analyze_docx(data)
    if suffix == ".xlsx":
        return _analyze_xlsx(data)
    if suffix == ".pptx":
        return _analyze_pptx(data)
    raise ValueError("Formato no soportado.")


def _money(value: float) -> str:
    return f"${value:,.4f}"


def render_print_cost_analyzer() -> None:
    st.title("Análisis y costeo de impresión")
    st.caption("Analiza cobertura CMYK y estima papel, tinta, energía, tiempo, depreciación, mantenimiento y desgaste.")

    uploaded = st.file_uploader("Archivo", type=SUPPORTED, help=f"Máximo recomendado: {MAX_FILE_MB} MB")
    if uploaded is None:
        st.info("Carga un PDF, JPG, PNG, DOCX, XLSX o PPTX para comenzar.")
        return
    if uploaded.size > MAX_FILE_MB * 1024 * 1024:
        st.error(f"El archivo supera {MAX_FILE_MB} MB.")
        return

    with st.expander("Configuración del trabajo", expanded=True):
        a, b, c = st.columns(3)
        copies = a.number_input("Copias", 1, 10000, 1)
        sides = b.selectbox("Impresión", ["Una cara", "Doble cara"])
        paper_name = c.selectbox("Papel", ["Bond carta", "Bond oficio", "Fotográfico mate", "Fotográfico brillante", "Opalina", "Adhesivo", "Personalizado"])
        a, b, c, d = st.columns(4)
        paper_cost = a.number_input("Costo por hoja ($)", min_value=0.0, value=0.0180, format="%.4f")
        waste_pct = b.number_input("Merma (%)", 0.0, 50.0, 3.0)
        margin_pct = c.number_input("Margen deseado (%)", 0.0, 500.0, 40.0)
        tax_pct = d.number_input("Impuestos/recargos (%)", 0.0, 100.0, 0.0)

    with st.expander("Impresora, tintas y costos técnicos"):
        profile = st.selectbox("Perfil", ["HP Smart Tank 580", "Epson EcoTank L3250", "Personalizado"])
        defaults = {
            "HP Smart Tank 580": (230.0, 50000, 70.0, 19.0, 6000, 12000, 10.0, 0.18),
            "Epson EcoTank L3250": (220.0, 50000, 65.0, 16.0, 7500, 4500, 10.0, 0.16),
            "Personalizado": (200.0, 40000, 70.0, 18.0, 6000, 10000, 8.0, 0.20),
        }[profile]
        a, b, c, d = st.columns(4)
        printer_cost = a.number_input("Costo impresora ($)", 0.0, value=defaults[0])
        life_pages = b.number_input("Vida útil estimada (páginas)", 1, value=defaults[1])
        bottle_ml = c.number_input("Contenido por tinta (ml)", 1.0, value=defaults[2])
        bottle_cost = d.number_input("Costo por botella ($)", 0.0, value=defaults[3])
        a, b, c, d = st.columns(4)
        color_yield = a.number_input("Rendimiento color al 5%", 1, value=defaults[4])
        black_yield = b.number_input("Rendimiento negro al 5%", 1, value=defaults[5])
        ppm = c.number_input("Velocidad real (pág/min)", 0.1, value=defaults[6])
        watts = d.number_input("Consumo imprimiendo (W)", 0.1, value=defaults[7] * 100)
        a, b, c = st.columns(3)
        electricity_kwh = a.number_input("Electricidad ($/kWh)", 0.0, value=0.10, format="%.4f")
        maintenance_per_page = b.number_input("Mantenimiento por página ($)", 0.0, value=0.0030, format="%.4f")
        labor_hour = c.number_input("Mano de obra por hora ($)", 0.0, value=2.50, format="%.2f")

    if not st.button("Analizar y calcular", type="primary"):
        return

    try:
        with st.spinner("Analizando cobertura y estructura del archivo..."):
            coverage = analyze_file(uploaded.name, uploaded.getvalue())
    except Exception as exc:
        st.error(f"No se pudo analizar el archivo: {exc}")
        return

    duplex_factor = .5 if sides == "Doble cara" else 1.0
    sheets = math.ceil(coverage.pages * copies * duplex_factor)
    billed_sheets = sheets * (1 + waste_pct / 100)
    printed_pages = coverage.pages * copies

    # Los rendimientos publicados se expresan normalmente a 5% de cobertura.
    ink_costs = {}
    for label, pct, yield_pages in (
        ("C", coverage.cyan, color_yield), ("M", coverage.magenta, color_yield),
        ("Y", coverage.yellow, color_yield), ("K", coverage.black, black_yield),
    ):
        equivalent_pages = printed_pages * (pct / 5.0)
        ink_costs[label] = equivalent_pages / yield_pages * bottle_cost

    paper_total = billed_sheets * paper_cost
    ink_total = sum(ink_costs.values())
    depreciation = printed_pages * printer_cost / life_pages
    maintenance = printed_pages * maintenance_per_page
    minutes = printed_pages / ppm + 2.0
    energy = (watts / 1000) * (minutes / 60) * electricity_kwh
    labor = minutes / 60 * labor_hour
    subtotal = paper_total + ink_total + depreciation + maintenance + energy + labor
    tax = subtotal * tax_pct / 100
    cost_total = subtotal + tax
    sale_price = cost_total / max(.01, 1 - margin_pct / 100) if margin_pct < 100 else cost_total * (1 + margin_pct / 100)

    st.subheader("Cobertura CMYK estimada")
    cols = st.columns(4)
    for col, label, value in zip(cols, ["Cian", "Magenta", "Amarillo", "Negro"], [coverage.cyan, coverage.magenta, coverage.yellow, coverage.black]):
        col.metric(label, f"{value:.1f}%")
    st.caption(f"Confianza: {coverage.confidence}. {coverage.note}")

    st.subheader("Resultado del costeo")
    cols = st.columns(5)
    cols[0].metric("Páginas", f"{printed_pages:,}")
    cols[1].metric("Hojas", f"{sheets:,}")
    cols[2].metric("Costo total", _money(cost_total))
    cols[3].metric("Costo/página", _money(cost_total / max(1, printed_pages)))
    cols[4].metric("Precio sugerido", f"${sale_price:,.2f}")

    rows = {
        "Papel y merma": paper_total,
        "Tinta C": ink_costs["C"], "Tinta M": ink_costs["M"], "Tinta Y": ink_costs["Y"], "Tinta K": ink_costs["K"],
        "Depreciación / vida útil": depreciation,
        "Mantenimiento y desgaste": maintenance,
        "Electricidad": energy,
        "Tiempo y mano de obra": labor,
        "Impuestos / recargos": tax,
    }
    st.dataframe([{"Concepto": key, "Costo ($)": round(value, 5), "% del costo": round(value / max(cost_total, .00001) * 100, 2)} for key, value in rows.items()], use_container_width=True, hide_index=True)

    st.warning("La cobertura CMYK es una estimación técnica. El consumo real depende del controlador, calidad, perfil ICC, limpieza de cabezales, sangrado, absorción del papel y estado de la impresora.")
